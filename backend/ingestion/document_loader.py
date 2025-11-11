"""
Document Loading and Processing
Supports multiple file formats for ingestion
"""
import os
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from abc import ABC, abstractmethod

logger = logging.getLogger(__name__)


class Document:
    """Represents a loaded document"""
    def __init__(self, content: str, metadata: Dict[str, Any], source_path: Path):
        self.content = content
        self.metadata = metadata
        self.source_path = source_path
        self.id = self._generate_id()

    def _generate_id(self) -> str:
        """Generate unique document ID"""
        # Use file path relative to data directory
        try:
            from ..config import settings
            relative_path = self.source_path.relative_to(settings.data_dir / "documents")
            return str(relative_path).replace(os.sep, '/')
        except ValueError:
            # Fallback to full path hash
            return str(hash(str(self.source_path)))

    def get_chunks(self, chunk_size: int = 512, overlap: int = 50) -> List[Dict[str, Any]]:
        """
        Split document into chunks

        Args:
            chunk_size: Maximum characters per chunk
            overlap: Characters to overlap between chunks

        Returns:
            List of chunk dictionaries with text and metadata
        """
        if not self.content:
            return []

        chunks = []
        start = 0

        while start < len(self.content):
            end = start + chunk_size

            # Try to break at sentence boundaries
            if end < len(self.content):
                # Look for sentence endings within the last 100 characters
                search_start = max(start, end - 100)
                sentence_end = self.content.rfind('.', search_start, end)
                if sentence_end == -1:
                    sentence_end = self.content.rfind('!', search_start, end)
                if sentence_end == -1:
                    sentence_end = self.content.rfind('?', search_start, end)
                if sentence_end == -1:
                    sentence_end = self.content.rfind('\n', search_start, end)

                if sentence_end != -1 and sentence_end > start:
                    end = sentence_end + 1

            chunk_text = self.content[start:end].strip()

            if chunk_text:  # Only add non-empty chunks
                chunk_metadata = self.metadata.copy()
                chunk_metadata.update({
                    'chunk_start': start,
                    'chunk_end': end,
                    'chunk_index': len(chunks),
                    'document_id': self.id
                })

                chunks.append({
                    'text': chunk_text,
                    'metadata': chunk_metadata
                })

            start = end - overlap

        return chunks


class DocumentLoader(ABC):
    """Abstract base class for document loaders"""

    @abstractmethod
    def can_load(self, file_path: Path) -> bool:
        """Check if this loader can handle the file"""
        pass

    @abstractmethod
    def load(self, file_path: Path) -> Optional[Document]:
        """Load document from file"""
        pass


class TextLoader(DocumentLoader):
    """Loader for plain text files"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'TextLoader'
            }

            return Document(content, metadata, file_path)

        except Exception as e:
            logger.error(f"Failed to load text file {file_path}: {e}")
            return None


class PDFLoader(DocumentLoader):
    """Loader for PDF files"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pdf'

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            content = ""

            for page in reader.pages:
                content += page.extract_text() + "\n"

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'PDFLoader',
                'num_pages': len(reader.pages)
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("PyPDF2 not installed, PDF loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load PDF file {file_path}: {e}")
            return None


class WordLoader(DocumentLoader):
    """Loader for Word documents"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.docx', '.doc']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            import docx

            doc = docx.Document(file_path)
            content = ""

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content += " | ".join(row_text) + "\n"
                content += "\n"

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'WordLoader',
                'num_paragraphs': len(doc.paragraphs),
                'num_tables': len(doc.tables)
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("python-docx not installed, Word document loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load Word file {file_path}: {e}")
            return None


class ExcelLoader(DocumentLoader):
    """Loader for Excel spreadsheets"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.xlsx', '.xls']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            import pandas as pd

            # Read all sheets
            excel_data = pd.read_excel(file_path, sheet_name=None)

            content = ""
            sheet_info = []

            for sheet_name, df in excel_data.items():
                content += f"\n--- Sheet: {sheet_name} ---\n"

                # Add column headers
                if not df.empty:
                    content += "Columns: " + ", ".join(df.columns.astype(str)) + "\n"

                    # Add sample data (first 10 rows)
                    sample_data = df.head(10)
                    content += sample_data.to_string(index=False) + "\n"

                    # Add summary statistics for numeric columns
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    if len(numeric_cols) > 0:
                        content += "\nSummary Statistics:\n"
                        content += df[numeric_cols].describe().to_string() + "\n"

                sheet_info.append({
                    'name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns)
                })

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'ExcelLoader',
                'sheets': sheet_info,
                'total_sheets': len(excel_data)
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("pandas/openpyxl/xlrd not installed, Excel loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load Excel file {file_path}: {e}")
            return None


class PowerPointLoader(DocumentLoader):
    """Loader for PowerPoint presentations"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.pptx', '.ppt']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content = ""

            for i, slide in enumerate(prs.slides):
                content += f"\n--- Slide {i+1} ---\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n"

                # Extract text from notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        content += f"Notes: {notes}\n"

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'PowerPointLoader',
                'num_slides': len(prs.slides)
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("python-pptx not installed, PowerPoint loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load PowerPoint file {file_path}: {e}")
            return None


class CSVLoader(DocumentLoader):
    """Loader for CSV and TSV files"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.csv', '.tsv']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            import pandas as pd

            # Auto-detect separator
            separators = [',', '\t', ';', '|']
            df = None
            separator = ','

            for sep in separators:
                try:
                    df = pd.read_csv(file_path, sep=sep, nrows=5)  # Read first 5 rows to test
                    if len(df.columns) > 1:  # If we get multiple columns, likely correct separator
                        separator = sep
                        break
                except:
                    continue

            # Read the full file with detected separator
            df = pd.read_csv(file_path, sep=separator)

            content = f"CSV File Analysis:\n"
            content += f"Separator: '{separator}'\n"
            content += f"Rows: {len(df)}\n"
            content += f"Columns: {len(df.columns)}\n\n"

            content += "Column Names: " + ", ".join(df.columns.astype(str)) + "\n\n"

            # Data types
            content += "Data Types:\n"
            for col, dtype in df.dtypes.items():
                content += f"  {col}: {dtype}\n"
            content += "\n"

            # Sample data
            content += "Sample Data (first 10 rows):\n"
            content += df.head(10).to_string(index=False) + "\n\n"

            # Summary for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                content += "Numeric Column Statistics:\n"
                content += df[numeric_cols].describe().to_string() + "\n"

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'CSVLoader',
                'rows': len(df),
                'columns': len(df.columns),
                'separator': separator
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("pandas not installed, CSV loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load CSV file {file_path}: {e}")
            return None


class YAMLLoader(DocumentLoader):
    """Loader for YAML configuration files"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.yaml', '.yml']

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            import yaml

            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)

            content = f"YAML Configuration File:\n\n"

            # Convert to readable text format
            def flatten_dict(d, prefix=''):
                result = ""
                if isinstance(d, dict):
                    for key, value in d.items():
                        full_key = f"{prefix}.{key}" if prefix else key
                        if isinstance(value, (dict, list)):
                            result += f"{full_key}:\n{flatten_dict(value, full_key)}\n"
                        else:
                            result += f"{full_key}: {value}\n"
                elif isinstance(d, list):
                    for i, item in enumerate(d):
                        full_key = f"{prefix}[{i}]"
                        if isinstance(item, (dict, list)):
                            result += f"{full_key}:\n{flatten_dict(item, full_key)}\n"
                        else:
                            result += f"{full_key}: {item}\n"
                else:
                    result += f"{prefix}: {d}\n"

                return result

            content += flatten_dict(data)

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'YAMLLoader'
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("PyYAML not installed, YAML loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load YAML file {file_path}: {e}")
            return None


class EPUBLoader(DocumentLoader):
    """Loader for EPUB e-books"""

    def can_load(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.epub'

    def load(self, file_path: Path) -> Optional[Document]:
        try:
            import ebooklib
            from ebooklib import epub

            book = epub.read_epub(file_path)
            content = ""

            # Extract title and metadata
            title = book.get_metadata('DC', 'title')
            if title:
                content += f"Title: {title[0][0]}\n\n"

            author = book.get_metadata('DC', 'creator')
            if author:
                content += f"Author: {author[0][0]}\n\n"

            # Extract text from chapters
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    chapter_content = item.get_content().decode('utf-8')

                    # Simple HTML tag removal (basic)
                    import re
                    clean_text = re.sub(r'<[^>]+>', '', chapter_content)
                    clean_text = re.sub(r'\s+', ' ', clean_text).strip()

                    if clean_text:
                        content += clean_text + "\n\n"

            metadata = {
                'file_name': file_path.name,
                'file_size': file_path.stat().st_size,
                'file_type': file_path.suffix,
                'last_modified': file_path.stat().st_mtime,
                'loader': 'EPUBLoader',
                'title': title[0][0] if title else 'Unknown'
            }

            return Document(content, metadata, file_path)

        except ImportError:
            logger.warning("ebooklib not installed, EPUB loading disabled")
            return None
        except Exception as e:
            logger.error(f"Failed to load EPUB file {file_path}: {e}")
            return None


class DocumentProcessor:
    """Main document processing class"""

    def __init__(self):
        self.loaders = [
            TextLoader(),
            PDFLoader(),
            WordLoader(),
            ExcelLoader(),
            PowerPointLoader(),
            CSVLoader(),
            YAMLLoader(),
            EPUBLoader()
        ]

    def load_document(self, file_path: Path) -> Optional[Document]:
        """
        Load a document using the appropriate loader

        Args:
            file_path: Path to the document file

        Returns:
            Loaded Document object or None if loading failed
        """
        if not file_path.exists():
            logger.error(f"File does not exist: {file_path}")
            return None

        # Find appropriate loader
        for loader in self.loaders:
            if loader.can_load(file_path):
                logger.info(f"Loading {file_path} with {loader.__class__.__name__}")
                return loader.load(file_path)

        logger.warning(f"No suitable loader found for {file_path}")
        return None

    def load_directory(self, directory_path: Path, recursive: bool = True) -> List[Document]:
        """
        Load all documents from a directory

        Args:
            directory_path: Path to directory
            recursive: Whether to search subdirectories

        Returns:
            List of loaded documents
        """
        if not directory_path.exists() or not directory_path.is_dir():
            logger.error(f"Directory does not exist: {directory_path}")
            return []

        documents = []
        pattern = "**/*" if recursive else "*"

        for file_path in directory_path.glob(pattern):
            if file_path.is_file():
                doc = self.load_document(file_path)
                if doc:
                    documents.append(doc)

        logger.info(f"Loaded {len(documents)} documents from {directory_path}")
        return documents

    def process_documents(self, documents: List[Document],
                         chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
        """
        Process documents into chunks

        Args:
            documents: List of Document objects
            chunk_size: Size of text chunks (uses config default if None)
            chunk_overlap: Overlap between chunks (uses config default if None)

        Returns:
            List of chunk dictionaries
        """
        from ..config import settings

        chunk_size = chunk_size or settings.max_chunk_size
        chunk_overlap = chunk_overlap or settings.chunk_overlap

        all_chunks = []

        for doc in documents:
            chunks = doc.get_chunks(chunk_size, chunk_overlap)
            all_chunks.extend(chunks)

        logger.info(f"Processed {len(documents)} documents into {len(all_chunks)} chunks")
        return all_chunks
